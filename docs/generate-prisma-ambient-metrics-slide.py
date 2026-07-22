#!/usr/bin/env python3
"""Generate Prisma Health ambient listening metrics slide (CipherHealth 2026 brand)."""

from __future__ import annotations

import argparse
import csv
import json
import sys
from collections import defaultdict
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

from render_utils import render_logo_png
from prisma_ambient_metrics import REPORT_META_KEY, hospital_metric_keys, load_metrics, pct

# CipherHealth 2026 Marketing Brand Colors
DARK_ANCHOR = RGBColor(0x2D, 0x26, 0x60)
DEEP_VIOLET = RGBColor(0x5B, 0x4F, 0xB5)
CIPHER_TEAL = RGBColor(0x00, 0xB4, 0xA4)
ELECTRIC_CHARTREUSE = RGBColor(0xB5, 0xCC, 0x2E)
WARM_AMBER = RGBColor(0xE8, 0xA0, 0x20)
LAVENDER_WHITE = RGBColor(0xF2, 0xF1, 0xF9)
TEAL_TINT = RGBColor(0xD0, 0xF4, 0xF1)
TEAL_TINT_TEXT = RGBColor(0x00, 0x4E, 0x48)
CHARTREUSE_TINT = RGBColor(0xE8, 0xF4, 0xA8)
CHARTREUSE_TEXT = RGBColor(0x2A, 0x32, 0x00)
AMBER_TINT = RGBColor(0xFD, 0xEF, 0xD0)
AMBER_TEXT = RGBColor(0x5C, 0x38, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_200 = RGBColor(0xE5, 0xE4, 0xED)

ROOT = Path(__file__).resolve().parent.parent
OUT = Path(__file__).resolve().parent / "prisma-health-ambient-listening-metrics-2026.pptx"
CSV_DEFAULT = Path("/Users/mcarroll/Desktop/BP Scripts/PrismaALInteractions.csv")
LOGO_SVG = ROOT / "assets" / "cipherhealth-logo-mark.svg"
LOGO_PNG = Path(__file__).resolve().parent / ".cache" / "cipherhealth-logo-mark.png"

DATE_START = datetime(2026, 5, 5)
DATE_END = datetime(2026, 7, 21, 23, 59, 59)
BAPTIST_CUTOFF = datetime(2026, 7, 1)

HOSPITALS_PPTX = [
    ("GREENVILLE MEMORIAL HOSPITAL", "Greenville Memorial Hospital", "GMH ", CIPHER_TEAL, TEAL_TINT, TEAL_TINT_TEXT),
    ("PATEWOOD MEMORIAL HOSPITAL", "Patewood Memorial Hospital", "PWH ", DEEP_VIOLET, LAVENDER_WHITE, DARK_ANCHOR),
    ("BAPTIST COLUMBIA HOSPITAL", "Baptist Columbia Hospital", "BCH ", WARM_AMBER, AMBER_TINT, AMBER_TEXT),
]

MARGIN_X = Inches(0.38)
HEADER_H = Inches(0.62)
ACCENT_H = Inches(0.035)
LOGO_SIZE = Inches(0.38)
TITLE_TOP = Inches(0.72)
TITLE_H = Inches(0.62)
CONTENT_TOP = Inches(1.42)
FOOTER_TOP = Inches(6.72)
FOOTER_H = Inches(0.48)


def ensure_logo_png() -> Path:
    if LOGO_PNG.exists() and LOGO_PNG.stat().st_mtime >= LOGO_SVG.stat().st_mtime:
        return LOGO_PNG
    render_logo_png(LOGO_SVG, LOGO_PNG, size=96)
    return LOGO_PNG


def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_textbox(slide, left, top, width, height, text, size=9, bold=False, color=DARK_ANCHOR, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_header(slide, prs, logo_path: Path, customer_name: str = ""):
    add_rect(slide, 0, 0, prs.slide_width, HEADER_H, DARK_ANCHOR)
    add_rect(slide, 0, HEADER_H - Pt(1), prs.slide_width, ACCENT_H, CIPHER_TEAL)

    logo_top = (HEADER_H - LOGO_SIZE) / 2
    slide.shapes.add_picture(str(logo_path), MARGIN_X, logo_top, width=LOGO_SIZE, height=LOGO_SIZE)

    add_textbox(
        slide,
        MARGIN_X + LOGO_SIZE + Inches(0.1),
        Inches(0.16),
        Inches(1.8),
        Inches(0.3),
        "CipherHealth",
        size=13,
        bold=True,
        color=WHITE,
    )

    if customer_name:
        add_textbox(
            slide,
            Inches(9.2),
            Inches(0.12),
            Inches(3.7),
            Inches(0.28),
            customer_name,
            size=11,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.RIGHT,
        )
        add_textbox(
            slide,
            Inches(9.2),
            Inches(0.34),
            Inches(3.7),
            Inches(0.2),
            "CipherRounds · Ambient Listening · 2026",
            size=7,
            bold=True,
            color=CIPHER_TEAL,
            align=PP_ALIGN.RIGHT,
        )


def format_provider_label(care_provider: str | None) -> str:
    return (care_provider or "").strip()


def title_with_provider(base_title: str, provider_label: str) -> str:
    if not provider_label:
        return base_title
    return f"{base_title} · {provider_label}"


def add_title_band(slide, prs, date_range_label: str, provider_label: str = ""):
    add_rect(slide, 0, HEADER_H, prs.slide_width, TITLE_H, LAVENDER_WHITE)
    add_textbox(
        slide,
        MARGIN_X,
        TITLE_TOP,
        Inches(8.5),
        Inches(0.28),
        "Ambient Listening Adoption Metrics",
        size=16,
        bold=True,
        color=DARK_ANCHOR,
    )
    subtitle = date_range_label
    if provider_label:
        subtitle = f"{date_range_label}  ·  {provider_label}"
    add_textbox(
        slide,
        MARGIN_X,
        TITLE_TOP + Inches(0.28),
        Inches(11.8),
        Inches(0.22),
        subtitle,
        size=8,
        color=GRAY_500,
    )


def add_metric_pill(slide, left, top, width, label, value, accent, bg, text_color):
    card = add_rect(slide, left, top, width, Inches(0.58), bg, accent)
    card.line.width = Pt(1)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.06), width - Inches(0.16), Inches(0.14), label, size=6.5, bold=True, color=text_color)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.24), width - Inches(0.16), Inches(0.24), value, size=14, bold=True, color=accent)


def add_hospital_column(slide, left, width, hospital_name, metrics, accent, bg, text_color):
    top = CONTENT_TOP
    height = FOOTER_TOP - CONTENT_TOP - Inches(0.08)

    card = add_rect(slide, left, top, width, height, bg, accent)
    card.line.width = Pt(1.5)
    add_rect(slide, left, top, width, Inches(0.05), accent)

    inner_left = left + Inches(0.1)
    inner_width = width - Inches(0.2)
    y = top + Inches(0.12)

    short_name = hospital_name.replace(" Hospital", "")
    add_textbox(slide, inner_left, y, inner_width, Inches(0.2), short_name, size=10, bold=True, color=DARK_ANCHOR)
    y += Inches(0.2)
    add_textbox(
        slide,
        inner_left,
        y,
        inner_width,
        Inches(0.14),
        f"{metrics['total_rounds']:,} total rounds",
        size=7,
        color=GRAY_500,
    )
    y += Inches(0.22)

    pill_w = (inner_width - Inches(0.08)) / 3
    add_metric_pill(
        slide,
        inner_left,
        y,
        pill_w,
        "AL Used",
        f"{metrics['al_used_pct']}%",
        accent,
        WHITE,
        text_color,
    )
    add_metric_pill(
        slide,
        inner_left + pill_w + Inches(0.04),
        y,
        pill_w,
        "Patient Said No",
        f"{metrics['patient_said_no_pct']}%",
        accent,
        WHITE,
        text_color,
    )
    add_metric_pill(
        slide,
        inner_left + (pill_w + Inches(0.04)) * 2,
        y,
        pill_w,
        "Manual Round",
        f"{metrics['manual_round_pct']}%",
        accent,
        WHITE,
        text_color,
    )
    y += Inches(0.68)

    add_textbox(
        slide,
        inner_left,
        y,
        inner_width,
        Inches(0.14),
        "ROUND COUNTS",
        size=6.5,
        bold=True,
        color=DEEP_VIOLET,
    )
    y += Inches(0.16)

    counts = [
        ("AL used", metrics["al_used"]),
        ("Patient said no", metrics["patient_said_no"]),
        ("Manual round", metrics["manual_round"]),
    ]
    if metrics["declined_during_al"]:
        counts.append(("Declined during AL", metrics["declined_during_al"]))
    for label, count in counts:
        add_textbox(slide, inner_left, y, inner_width * 0.68, Inches(0.13), label, size=6, color=DARK_ANCHOR)
        add_textbox(
            slide,
            inner_left + inner_width * 0.68,
            y,
            inner_width * 0.32,
            Inches(0.13),
            f"{count:,}",
            size=6,
            bold=True,
            color=text_color,
            align=PP_ALIGN.RIGHT,
        )
        y += Inches(0.125)

    if metrics["declined_or_no"]:
        y += Inches(0.04)
        add_rect(slide, inner_left, y, inner_width, Inches(0.42), WHITE, accent)
        add_textbox(
            slide,
            inner_left + Inches(0.06),
            y + Inches(0.04),
            inner_width - Inches(0.12),
            Inches(0.14),
            "PATIENT/STAFF NO TO AL",
            size=5.5,
            bold=True,
            color=DEEP_VIOLET,
        )
        add_textbox(
            slide,
            inner_left + Inches(0.06),
            y + Inches(0.18),
            inner_width * 0.4,
            Inches(0.18),
            f"{metrics['declined_or_no_pct']}%",
            size=14,
            bold=True,
            color=accent,
        )
        add_textbox(
            slide,
            inner_left + Inches(0.06),
            y + Inches(0.3),
            inner_width - Inches(0.12),
            Inches(0.1),
            f"{metrics['declined_or_no']:,} of {metrics['total_rounds']:,}",
            size=5.5,
            color=GRAY_500,
        )
        y += Inches(0.46)

    units_with_al = sum(1 for u in metrics["units"] if u["al_used"] + u["declined_during_al"] > 0)
    add_textbox(
        slide,
        inner_left,
        top + height - Inches(0.22),
        inner_width,
        Inches(0.16),
        f"{len(metrics['units'])} units · {units_with_al} with AL",
        size=6,
        color=GRAY_500,
    )


UNIT_TABLE_TOP = Inches(1.38)
UNIT_TABLE_BOTTOM = Inches(6.62)


def add_unit_detail_title(slide, prs, metrics, date_range_label: str, provider_label: str = ""):
    add_rect(slide, 0, HEADER_H, prs.slide_width, Inches(0.56), LAVENDER_WHITE)
    add_textbox(
        slide,
        MARGIN_X,
        Inches(0.78),
        Inches(11),
        Inches(0.24),
        title_with_provider("Unit / Department Breakdown", provider_label),
        size=14,
        bold=True,
        color=DARK_ANCHOR,
    )
    add_textbox(
        slide,
        MARGIN_X,
        Inches(1.02),
        Inches(11.8),
        Inches(0.16),
        f"{date_range_label}  ·  {metrics['total_rounds']:,} total rounds",
        size=7.5,
        color=GRAY_500,
    )
    add_textbox(
        slide,
        MARGIN_X,
        Inches(1.18),
        Inches(11.8),
        Inches(0.16),
        f"{metrics['al_used']:,} AL used  ·  {metrics['patient_said_no']:,} patient said no  ·  "
        f"{metrics['manual_round']:,} manual  ·  {metrics['declined_during_al']:,} declined during AL",
        size=7,
        color=GRAY_500,
    )


def add_unit_table(slide, prs, units: list[dict], accent, prefix: str):
    table_left = MARGIN_X
    table_width = prs.slide_width - MARGIN_X * 2
    y = UNIT_TABLE_TOP

    headers = [
        ("Unit / Department", 0.20, PP_ALIGN.LEFT),
        ("Total\nRounds", 0.08, PP_ALIGN.CENTER),
        ("AL\nUsed", 0.08, PP_ALIGN.CENTER),
        ("Patient\nSaid No", 0.08, PP_ALIGN.CENTER),
        ("Manual\nRound", 0.08, PP_ALIGN.CENTER),
        ("Declined\nDuring AL", 0.08, PP_ALIGN.CENTER),
        ("% Patient/Staff\nNo to AL", 0.09, PP_ALIGN.CENTER),
    ]

    header_h = Inches(0.32)
    add_rect(slide, table_left, y, table_width, header_h, accent)
    x = table_left
    for label, frac, align in headers:
        col_w = table_width * frac
        add_textbox(
            slide,
            x + Inches(0.04),
            y + Inches(0.02),
            col_w - Inches(0.06),
            header_h,
            label,
            size=6,
            bold=True,
            color=WHITE,
            align=align,
        )
        x += col_w
    y += header_h

    row_h = Inches(0.132)
    max_rows = int((UNIT_TABLE_BOTTOM - y) / row_h)
    display_units = units[:max_rows]

    for idx, unit in enumerate(display_units):
        row_bg = WHITE if idx % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFC)
        add_rect(slide, table_left, y, table_width, row_h, row_bg, GRAY_200)
        x = table_left
        unit_label = unit["unit"].replace(prefix, "")
        values = [
            (unit_label, headers[0][1], DARK_ANCHOR, PP_ALIGN.LEFT, False),
            (f"{unit['total']:,}", headers[1][1], DARK_ANCHOR, PP_ALIGN.CENTER, False),
            (str(unit["al_used"]), headers[2][1], DARK_ANCHOR, PP_ALIGN.CENTER, unit["al_used"] > 0),
            (str(unit["patient_said_no"]), headers[3][1], DARK_ANCHOR, PP_ALIGN.CENTER, unit["patient_said_no"] > 0),
            (f"{unit['manual_round']:,}", headers[4][1], GRAY_500, PP_ALIGN.CENTER, unit["manual_round"] > 0),
            (
                str(unit["declined_during_al"]),
                headers[5][1],
                DARK_ANCHOR,
                PP_ALIGN.CENTER,
                unit["declined_during_al"] > 0,
            ),
            (
                f"{unit['declined_or_no_pct']}%",
                headers[6][1],
                DARK_ANCHOR,
                PP_ALIGN.CENTER,
                unit["declined_or_no"] > 0,
            ),
        ]
        for text, frac, color, align, bold in values:
            col_w = table_width * frac
            add_textbox(
                slide,
                x + Inches(0.04),
                y + Inches(0.015),
                col_w - Inches(0.06),
                row_h,
                text,
                size=5.8,
                bold=bold,
                color=color,
                align=align,
            )
            x += col_w
        y += row_h

    if len(units) > len(display_units):
        add_textbox(
            slide,
            table_left,
            y + Inches(0.04),
            table_width,
            Inches(0.16),
            f"+ {len(units) - len(display_units)} additional units not shown",
            size=6.5,
            color=GRAY_500,
        )


def add_baptist_after_july_slide(
    prs,
    logo_path: Path,
    comparison: dict,
    prefix: str,
    slide_num: int,
    slide_total: int,
    provider_label: str = "",
    customer_name: str = "",
):
    after = comparison["after"]
    period_label = comparison["after_period"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, prs, logo_path, customer_name)

    add_rect(slide, 0, HEADER_H, prs.slide_width, Inches(0.56), LAVENDER_WHITE)
    add_textbox(
        slide,
        MARGIN_X,
        Inches(0.78),
        Inches(11),
        Inches(0.24),
        title_with_provider("Baptist Columbia · After July 1", provider_label),
        size=14,
        bold=True,
        color=DARK_ANCHOR,
    )
    add_textbox(
        slide,
        MARGIN_X,
        Inches(1.02),
        Inches(11.8),
        Inches(0.16),
        f"{period_label}  ·  {after['total_rounds']:,} total rounds",
        size=7.5,
        color=GRAY_500,
    )
    add_textbox(
        slide,
        MARGIN_X,
        Inches(1.18),
        Inches(11.8),
        Inches(0.16),
        f"{after['al_used']:,} AL used ({after['al_used_pct']}%)  ·  "
        f"{after['patient_said_no']:,} patient said no  ·  "
        f"{after['manual_round']:,} manual  ·  "
        f"{after['declined_during_al']:,} declined during AL",
        size=7,
        color=GRAY_500,
    )

    add_unit_table(slide, prs, after["units"], WARM_AMBER, prefix)
    add_footer(slide, prs, slide_num, slide_total)


def add_unit_detail_slide(
    prs,
    logo_path: Path,
    hospital_name,
    metrics,
    accent,
    prefix: str,
    slide_num: int,
    slide_total: int,
    date_range_label: str,
    provider_label: str = "",
    customer_name: str = "",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, prs, logo_path, customer_name)
    add_unit_detail_title(slide, prs, metrics, date_range_label, provider_label)
    add_unit_table(slide, prs, metrics["units"], accent, prefix)
    add_footer(slide, prs, slide_num, slide_total)


def add_summary_band(slide, prs, metrics: dict, date_range_label: str):
    band_top = FOOTER_TOP - Inches(0.72)
    add_rect(slide, MARGIN_X, band_top, prs.slide_width - (MARGIN_X * 2), Inches(0.56), CHARTREUSE_TINT, ELECTRIC_CHARTREUSE)
    add_textbox(slide, MARGIN_X + Inches(0.12), band_top + Inches(0.08), Inches(1.5), Inches(0.12), "KEY INSIGHTS", size=6.5, bold=True, color=CHARTREUSE_TEXT)
    hospital_metrics = {key: metrics[key] for key in hospital_metric_keys(metrics)}
    total_rounds = sum(item["total_rounds"] for item in hospital_metrics.values())
    total_used = sum(item["al_used"] for item in hospital_metrics.values())
    total_manual = sum(item["manual_round"] for item in hospital_metrics.values())
    total_during_al = sum(item["declined_during_al"] for item in hospital_metrics.values())
    used_pct = pct(total_used, total_rounds)
    manual_pct = pct(total_manual, total_rounds)
    during_pct = pct(total_during_al, total_rounds)
    baptist = metrics[REPORT_META_KEY]["baptist_before_after"]
    insights = [
        f"{date_range_label}: {total_rounds:,} rounds across 3 hospitals",
        f"AL used {used_pct}% ({total_used:,}); declined during AL {during_pct}% ({total_during_al:,}); manual {manual_pct}%",
        f"Baptist after July 1: {baptist['after']['total_rounds']:,} rounds, AL used {baptist['after']['al_used_pct']}% — see slide 5",
    ]
    x = MARGIN_X + Inches(1.55)
    col_w = (prs.slide_width - MARGIN_X * 2 - Inches(1.55)) / 3
    for idx, insight in enumerate(insights):
        add_textbox(slide, x + col_w * idx, band_top + Inches(0.08), col_w - Inches(0.08), Inches(0.48), insight, size=6.2, color=CHARTREUSE_TEXT)


def add_footer(slide, prs, slide_num: int | None = None, slide_total: int | None = None):
    add_rect(slide, 0, FOOTER_TOP, prs.slide_width, FOOTER_H, DARK_ANCHOR)
    add_textbox(
        slide,
        MARGIN_X,
        FOOTER_TOP + Inches(0.06),
        Inches(8.8),
        Inches(0.36),
        "Source: interaction export · % Pt/Staff No = Declined Yes ÷ total · AL Used = Amb Yes + Decl No · Manual = Amb No + Decl No",
        size=6,
        color=WHITE,
    )
    right_text = "Confidential · cipherhealth.com"
    if slide_num is not None and slide_total is not None:
        right_text = f"Slide {slide_num} of {slide_total} · {right_text}"
    add_textbox(
        slide,
        Inches(9.6),
        FOOTER_TOP + Inches(0.16),
        Inches(3.3),
        Inches(0.2),
        right_text,
        size=6.5,
        bold=True,
        color=CIPHER_TEAL,
        align=PP_ALIGN.RIGHT,
    )


def build_slide(
    csv_path: Path,
    care_provider: str | None = None,
    rounder: str | None = None,
    unit: str | None = None,
    output_path: Path | None = None,
) -> tuple[Path, dict]:
    logo_path = ensure_logo_png()
    metrics = load_metrics(csv_path, care_provider=care_provider, rounder=rounder, unit=unit)
    out_path = output_path or OUT
    report = metrics[REPORT_META_KEY]
    date_range_label = report["date_range_label"]
    provider_label = format_provider_label(care_provider)
    customer_name = provider_label

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_header(slide, prs, logo_path, customer_name)
    add_title_band(slide, prs, date_range_label, provider_label)

    usable = prs.slide_width - (MARGIN_X * 2)
    gap = Inches(0.12)
    col_w = (usable - gap * 2) / 3
    for idx, (_, hospital_name, _, accent, bg, text_color) in enumerate(HOSPITALS_PPTX):
        left = MARGIN_X + (col_w + gap) * idx
        add_hospital_column(slide, left, col_w, hospital_name, metrics[hospital_name], accent, bg, text_color)

    add_summary_band(slide, prs, metrics, date_range_label)
    slide_total = 1 + len(HOSPITALS_PPTX) + 1
    add_footer(slide, prs, 1, slide_total)

    for slide_idx, (_, hospital_name, prefix, accent, _, _) in enumerate(HOSPITALS_PPTX, start=2):
        add_unit_detail_slide(
            prs,
            logo_path,
            hospital_name,
            metrics[hospital_name],
            accent,
            prefix,
            slide_idx,
            slide_total,
            date_range_label,
            provider_label,
            customer_name,
        )

    baptist_prefix = HOSPITALS_PPTX[2][2]
    add_baptist_after_july_slide(
        prs,
        logo_path,
        report["baptist_before_after"],
        baptist_prefix,
        slide_total,
        slide_total,
        provider_label,
        customer_name,
    )

    prs.save(out_path)
    return out_path, metrics


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Prisma Health ambient listening metrics deck")
    parser.add_argument("csv", nargs="?", default=str(CSV_DEFAULT), help="Path to interaction export CSV")
    parser.add_argument(
        "--care-provider",
        dest="care_provider",
        help="Hospital / care provider name for deck branding (e.g. Prisma, Kaiser).",
    )
    parser.add_argument("--rounder", help="Filter to rounder (Interaction Created By).")
    parser.add_argument("--unit", help="Filter to unit.")
    parser.add_argument("-o", "--output", default=str(OUT), help="Output PowerPoint path")
    args = parser.parse_args()

    csv_path = Path(args.csv)
    out_path = Path(args.output)
    _, metrics = build_slide(
        csv_path,
        care_provider=args.care_provider,
        rounder=args.rounder,
        unit=args.unit,
        output_path=out_path,
    )
    cache_dir = Path(__file__).resolve().parent / ".cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    metrics_path = cache_dir / "prisma-ambient-metrics.json"
    metrics_path.write_text(json.dumps(metrics, indent=2))
    print(f"Wrote {out_path}")
    print(f"Metrics saved to {metrics_path}")


if __name__ == "__main__":
    main()
